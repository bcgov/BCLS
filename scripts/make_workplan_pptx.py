"""Generate a clean, professional Look West work-plan PowerPoint deck."""
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---- palette ----
NAVY = RGBColor(0x00, 0x2B, 0x57)
NAVY2 = RGBColor(0x0A, 0x3D, 0x75)
GOLD = RGBColor(0xFC, 0xBA, 0x19)
INK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
CARD = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD9, 0xE2, 0xEC)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
BLUE = RGBColor(0x25, 0x63, 0xEB)

FONT = "Segoe UI"
FONT_L = "Segoe UI Light"

EMUW, EMUH = Inches(13.333), Inches(7.5)


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _noline(shape):
    shape.line.fill.background()


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of dicts {text, size, color, bold, italic, space_after, font}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        if ln.get("space_after") is not None:
            p.space_after = Pt(ln["space_after"])
        if ln.get("space_before") is not None:
            p.space_before = Pt(ln["space_before"])
        if ln.get("bullet"):
            p.level = ln.get("level", 0)
        run = p.add_run()
        run.text = ln["text"]
        f = run.font
        f.size = Pt(ln.get("size", 18))
        f.bold = ln.get("bold", False)
        f.italic = ln.get("italic", False)
        f.name = ln.get("font", FONT)
        f.color.rgb = ln.get("color", INK)
    return tb


def card(slide, x, y, w, h, fill=CARD, line=LINE, line_w=Pt(1), radius=True, shadow=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w
    shp.shadow.inherit = False
    return shp


def accent_bar(slide, x, y, w, color=GOLD, h=Pt(5)):
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _solid(bar, color)
    bar.shadow.inherit = False
    return bar


def base(slide, prs, kicker, title, page=None):
    """Standard content-slide header. Returns y where body can start."""
    # background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMUW, EMUH)
    _solid(bg, WHITE)
    bg.shadow.inherit = False
    # side rail
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), EMUH)
    _solid(rail, NAVY)
    rail.shadow.inherit = False
    accent_bar(slide, Inches(0.55), Inches(0.62), Inches(0.7))
    textbox(slide, Inches(0.55), Inches(0.74), Inches(11), Inches(0.4),
            [{"text": kicker.upper(), "size": 12, "color": GOLD, "bold": True}])
    textbox(slide, Inches(0.55), Inches(1.05), Inches(12.2), Inches(0.9),
            [{"text": title, "size": 30, "color": NAVY, "bold": True, "font": FONT_L}])
    # footer
    textbox(slide, Inches(0.55), Inches(7.02), Inches(8), Inches(0.3),
            [{"text": "Look West Monitoring Dashboard & Data Workflow", "size": 9, "color": MUTED}])
    if page is not None:
        textbox(slide, Inches(12.2), Inches(7.02), Inches(0.8), Inches(0.3),
                [{"text": str(page), "size": 9, "color": MUTED, "align": PP_ALIGN.RIGHT}])
    return Inches(2.05)


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ===================================================================
def build(dst):
    prs = Presentation()
    prs.slide_width = EMUW
    prs.slide_height = EMUH

    # ---------- 1. TITLE ----------
    s = add_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMUW, EMUH)
    bg.fill.gradient()
    bg.fill.gradient_stops[0].color.rgb = NAVY
    bg.fill.gradient_stops[1].color.rgb = NAVY2
    try:
        bg.fill.gradient_angle = 30.0
    except Exception:
        pass
    _noline(bg)
    bg.shadow.inherit = False
    # gold accent block
    accent_bar(s, Inches(0.9), Inches(2.5), Inches(1.4), GOLD, Pt(6))
    textbox(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(0.5),
            [{"text": "LOOK WEST STRATEGY", "size": 16, "color": GOLD, "bold": True}])
    textbox(s, Inches(0.9), Inches(3.15), Inches(11.6), Inches(1.7),
            [{"text": "Monitoring Dashboard & Data Workflow", "size": 46, "color": WHITE,
              "bold": True, "font": FONT_L, "space_after": 6},
             {"text": "A two-track work plan: trusted data  →  a single source of truth  →  an executive-ready dashboard",
              "size": 18, "color": RGBColor(0xCB, 0xD9, 0xEC)}])
    # meta strip
    textbox(s, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.8),
            [{"text": "Work Plan  ·  Draft v0.2  ·  June 2026", "size": 13, "color": WHITE, "bold": True,
              "space_after": 2},
             {"text": "Owner: Mehdi Naji   ·   For review & sign-off: Jacqueline (Director)", "size": 12,
              "color": RGBColor(0xAF, 0xC2, 0xDB)}])

    # ---------- 2. AGENDA / AT A GLANCE ----------
    s = add_slide(prs)
    y = base(prs and prs, s, prs) if False else None
    y = base(s, prs, "Overview", "What this plan covers", 2)
    items = [
        ("01", "The challenge & approach", "Why monitoring needs one source of truth"),
        ("02", "Two parallel tracks", "Data hub  +  dashboard & deployment"),
        ("03", "Metrics & governance", "How metrics are proposed, validated, approved"),
        ("04", "Workflow & engagement", "How the tracks connect; who we engage"),
        ("05", "Timeline to end of summer", "Milestones, checkpoints, dependencies"),
        ("06", "Risks & next steps", "Bottlenecks, mitigations, immediate actions"),
    ]
    cols = 3
    cw, ch = Inches(3.95), Inches(1.95)
    gx, gy = Inches(0.2), Inches(0.3)
    x0, y0 = Inches(0.55), Inches(2.3)
    for i, (num, t, d) in enumerate(items):
        r, c = divmod(i, cols)
        x = x0 + c * (cw + gx)
        yy = y0 + r * (ch + gy)
        cd = card(s, x, yy, cw, ch, CARD, LINE)
        textbox(s, x + Inches(0.25), yy + Inches(0.2), Inches(2), Inches(0.6),
                [{"text": num, "size": 30, "color": GOLD, "bold": True, "font": FONT_L}])
        textbox(s, x + Inches(0.25), yy + Inches(0.85), cw - Inches(0.5), Inches(1.0),
                [{"text": t, "size": 16, "color": NAVY, "bold": True, "space_after": 3},
                 {"text": d, "size": 12, "color": MUTED}])

    # ---------- 3. CURRENT CONTEXT / WHERE WE ARE ----------
    s = add_slide(prs)
    base(s, prs, "Context", "Where we are right now", 3)
    # left: the gap
    lx, lw = Inches(0.55), Inches(5.9)
    card(s, lx, Inches(2.25), lw, Inches(4.4), CARD, LINE)
    textbox(s, lx + Inches(0.35), Inches(2.5), lw - Inches(0.7), Inches(0.5),
            [{"text": "THE CHALLENGE", "size": 13, "color": GOLD, "bold": True}])
    textbox(s, lx + Inches(0.35), Inches(2.95), lw - Inches(0.7), Inches(3.5),
            [{"text": "Look West progress data is fragmented", "size": 18, "color": NAVY, "bold": True,
              "space_after": 10},
             {"text": "•  Spread across public stats, internal spreadsheets, and announcements",
              "size": 14, "color": INK, "space_after": 8},
             {"text": "•  No single validated source of truth", "size": 14, "color": INK, "space_after": 8},
             {"text": "•  No executive-facing view of progress, gaps, or risks", "size": 14, "color": INK,
              "space_after": 8},
             {"text": "•  Metrics not yet defined or leadership-approved", "size": 14, "color": INK}])
    # right: what exists
    rx, rw = Inches(6.85), Inches(5.93)
    card(s, rx, Inches(2.25), rw, Inches(4.4), NAVY, None)
    textbox(s, rx + Inches(0.35), Inches(2.5), rw - Inches(0.7), Inches(0.5),
            [{"text": "ALREADY IN PLACE", "size": 13, "color": GOLD, "bold": True}])
    for i, (t, d) in enumerate([
        ("Data Hub schema designed", "13-table model: targets, trackers, announcements, reference data"),
        ("Dashboard prototype built", "HTML Monitor + 5 tracker pages (flexible design & UX)"),
        ("Governance scaffolding", "Ownership registry + dataset status & cadence"),
        ("Not yet done", "Metrics approval · validated data · deployment to users"),
    ]):
        yy = Inches(3.0) + i * Inches(0.86)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, rx + Inches(0.35), yy + Inches(0.06), Inches(0.16), Inches(0.16))
        _solid(dot, GOLD if i < 3 else RGBColor(0xF8, 0x71, 0x71))
        textbox(s, rx + Inches(0.65), yy, rw - Inches(1.0), Inches(0.8),
                [{"text": t, "size": 14, "color": WHITE, "bold": True, "space_after": 1},
                 {"text": d, "size": 11, "color": RGBColor(0xC0, 0xCE, 0xE0)}])

    # ---------- 4. APPROACH: HUB-FIRST ----------
    s = add_slide(prs)
    base(s, prs, "Approach", "One rule: the Data Hub is the source of truth", 4)
    textbox(s, Inches(0.55), Inches(1.95), Inches(12), Inches(0.6),
            [{"text": "Data flows in one direction. Nothing is typed straight into the dashboard — every figure traces to a validated row in the Hub.",
              "size": 15, "color": MUTED, "italic": True}])
    # flow chips
    chips = [
        ("Sources", "Public data · internal teams · announcements", LIGHT, NAVY, NAVY),
        ("Look West Data Hub", "Validated · owned · governed  —  SOURCE OF TRUTH", NAVY, WHITE, GOLD),
        ("Look West Monitor", "Executive dashboard · progress, gaps, risks", LIGHT, NAVY, NAVY),
        ("Executives & users", "Decision-ready, accessible", LIGHT, NAVY, NAVY),
    ]
    n = len(chips)
    cw = Inches(2.85)
    gap = Inches(0.42)
    total = n * cw + (n - 1) * gap
    x = Inches((13.333 - total / 914400) / 2 * 914400 / 914400)  # center
    x = Emu(int((EMUW - total) / 2))
    yy = Inches(3.3)
    ch = Inches(2.0)
    for i, (t, d, fill, txt, acc) in enumerate(chips):
        cx = Emu(int(x) + i * int(cw + gap))
        c = card(s, cx, yy, cw, ch, fill, None if fill == NAVY else LINE)
        accent_bar(s, cx + Inches(0.25), yy + Inches(0.28), Inches(0.55), acc)
        textbox(s, cx + Inches(0.25), yy + Inches(0.5), cw - Inches(0.5), ch - Inches(0.7),
                [{"text": t, "size": 16, "color": txt, "bold": True, "space_after": 6},
                 {"text": d, "size": 11.5, "color": (RGBColor(0xC8,0xD6,0xE8) if fill == NAVY else MUTED)}])
        if i < n - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON, cx + cw + Inches(0.04),
                                    yy + ch / 2 - Inches(0.22), Inches(0.34), Inches(0.44))
            _solid(ar, GOLD)
            ar.shadow.inherit = False
    textbox(s, Inches(0.55), Inches(5.9), Inches(12.2), Inches(0.9),
            [{"text": "Why this matters:  the long pole of the project is sourcing and validating data — not building the page.",
              "size": 15, "color": NAVY, "bold": True}])

    # ---------- 5. TWO TRACKS ----------
    s = add_slide(prs)
    base(s, prs, "Approach", "Two parallel tracks, one connection point", 5)
    tracks = [
        ("WORKSTREAM 1", "Data gathering, validation & Data Hub",
         ["Define a metric for each goal", "Source from public / internal / announcements",
          "Clean, standardize & validate with teams", "Land it in the Hub as source of truth"],
         "Critical path", NAVY),
        ("WORKSTREAM 2", "Dashboard, reporting & deployment",
         ["Connect the dashboard to the Hub", "Report progress, flag gaps & risks",
          "Add BC economic context snapshot", "Deploy securely with government IT"],
         "Follows the data", NAVY2),
    ]
    cw, ch = Inches(5.95), Inches(4.5)
    for i, (k, t, pts, tag, col) in enumerate(tracks):
        x = Inches(0.55) + i * (cw + Inches(0.35))
        yy = Inches(2.15)
        card(s, x, yy, cw, ch, CARD, LINE)
        head = card(s, x, yy, cw, Inches(1.15), col, None)
        head.adjustments[0] = 0.06
        textbox(s, x + Inches(0.35), yy + Inches(0.2), cw - Inches(0.7), Inches(0.9),
                [{"text": k, "size": 12, "color": GOLD, "bold": True, "space_after": 2},
                 {"text": t, "size": 18, "color": WHITE, "bold": True}])
        for j, p in enumerate(pts):
            py = yy + Inches(1.45) + j * Inches(0.62)
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.4), py + Inches(0.07), Inches(0.14), Inches(0.14))
            _solid(dot, GOLD)
            textbox(s, x + Inches(0.68), py, cw - Inches(1.0), Inches(0.6),
                    [{"text": p, "size": 14, "color": INK}])
        tagb = card(s, x + Inches(0.35), yy + ch - Inches(0.72), Inches(2.4), Inches(0.45),
                    LIGHT, None)
        tagb.adjustments[0] = 0.5
        textbox(s, x + Inches(0.35), yy + ch - Inches(0.69), Inches(2.4), Inches(0.4),
                [{"text": tag, "size": 12, "color": col, "bold": True, "align": PP_ALIGN.CENTER}],
                anchor=MSO_ANCHOR.MIDDLE)
    # connector note
    textbox(s, Inches(0.55), Inches(6.75), Inches(12), Inches(0.4),
            [{"text": "Both tracks meet at the Data Hub. Deployment with IT starts early, in parallel — not at the end.",
              "size": 13, "color": MUTED, "italic": True, "align": PP_ALIGN.CENTER}])

    # ---------- 6. METRICS PROCESS ----------
    s = add_slide(prs)
    base(s, prs, "Metrics", "How every metric is proposed & approved", 6)
    steps = [
        ("Draft", "One metric per goal, using a standard template", GOLD),
        ("Consult", "Internal teams confirm it's meaningful & feasible", BLUE),
        ("Revise", "Refine based on team feedback", BLUE),
        ("Approve", "Leadership signs off (Jacqueline)", GREEN),
        ("Populate", "Locked into the Data Hub", NAVY),
    ]
    n = len(steps)
    cw = Inches(2.3)
    gap = Inches(0.28)
    total = n * cw + (n - 1) * gap
    x0 = Emu(int((EMUW - total) / 2))
    yy = Inches(2.3)
    for i, (t, d, col) in enumerate(steps):
        cx = Emu(int(x0) + i * int(cw + gap))
        card(s, cx, yy, cw, Inches(1.9), CARD, LINE)
        num = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + cw/2 - Inches(0.33), yy + Inches(0.22),
                                 Inches(0.66), Inches(0.66))
        _solid(num, col)
        textbox(s, cx + cw/2 - Inches(0.33), yy + Inches(0.22), Inches(0.66), Inches(0.66),
                [{"text": str(i+1), "size": 22, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER}],
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, cx + Inches(0.15), yy + Inches(1.0), cw - Inches(0.3), Inches(0.85),
                [{"text": t, "size": 15, "color": NAVY, "bold": True, "align": PP_ALIGN.CENTER, "space_after": 2},
                 {"text": d, "size": 10.5, "color": MUTED, "align": PP_ALIGN.CENTER}])
        if i < n - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON, cx + cw + Inches(0.01),
                                    yy + Inches(0.72), Inches(0.26), Inches(0.4))
            _solid(ar, LINE)
            ar.shadow.inherit = False
    # what we capture per metric
    card(s, Inches(0.55), Inches(4.6), Inches(12.23), Inches(1.95), NAVY, None)
    textbox(s, Inches(0.9), Inches(4.8), Inches(11.5), Inches(0.4),
            [{"text": "FOR EACH PROPOSED METRIC WE DOCUMENT", "size": 13, "color": GOLD, "bold": True}])
    fields = ["What it measures", "Why it's relevant", "Does reliable data exist?",
              "Public or internal source?", "Update frequency", "Limitations, gaps & risks"]
    for i, fld in enumerate(fields):
        r, c = divmod(i, 3)
        fx = Inches(0.9) + c * Inches(3.95)
        fy = Inches(5.25) + r * Inches(0.55)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, fx, fy + Inches(0.07), Inches(0.14), Inches(0.14))
        _solid(dot, GOLD)
        textbox(s, fx + Inches(0.28), fy, Inches(3.6), Inches(0.5),
                [{"text": fld, "size": 13, "color": WHITE}])

    # ---------- 7. GOVERNANCE & ENGAGEMENT ----------
    s = add_slide(prs)
    base(s, prs, "Governance", "Checkpoints, approvals & who we engage", 7)
    # left: checkpoints
    lx, lw = Inches(0.55), Inches(6.4)
    textbox(s, lx, Inches(2.0), lw, Inches(0.4),
            [{"text": "FIVE CHECKPOINTS", "size": 13, "color": GOLD, "bold": True}])
    cps = [
        ("CP1", "Metrics consultation", "Internal teams confirm feasibility", BLUE),
        ("CP2", "Leadership metric sign-off", "Approved set locked into Hub · Jacqueline", GREEN),
        ("CP3", "Data validation", "Each dataset validated, marked approved", BLUE),
        ("CP4", "Deployment go / no-go", "Hosting, access & refresh · IT + Jacqueline", GOLD),
        ("CP5", "Executive review", "Dashboard meets needs · continue / scale", GREEN),
    ]
    for i, (cp, t, d, col) in enumerate(cps):
        yy = Inches(2.45) + i * Inches(0.82)
        card(s, lx, yy, lw, Inches(0.7), CARD, LINE)
        badge = card(s, lx + Inches(0.12), yy + Inches(0.12), Inches(0.85), Inches(0.46), col, None)
        badge.adjustments[0] = 0.3
        textbox(s, lx + Inches(0.12), yy + Inches(0.12), Inches(0.85), Inches(0.46),
                [{"text": cp, "size": 14, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER}],
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, lx + Inches(1.12), yy + Inches(0.08), lw - Inches(1.3), Inches(0.6),
                [{"text": t, "size": 13.5, "color": NAVY, "bold": True, "space_after": 0},
                 {"text": d, "size": 10.5, "color": MUTED}])
    # right: engagement
    rx, rw = Inches(7.25), Inches(5.55)
    textbox(s, rx, Inches(2.0), rw, Inches(0.4),
            [{"text": "WHO WE ENGAGE", "size": 13, "color": GOLD, "bold": True}])
    people = [
        ("Jacqueline", "Director — approval", "CP2 metrics · CP4 deploy · CP5 review"),
        ("Daniel", "Analyst — data input", "Sourcing, cleaning, validation into the Hub"),
        ("Nathan", "Technical / IT advice", "Hosting, security & deployment path (CP4)"),
        ("Internal teams / SMEs", "Validate & confirm", "CP1 consultation · data validation"),
        ("Mehdi", "Project lead", "Design, coordination, delivery, narrative"),
    ]
    for i, (nm, role, when) in enumerate(people):
        yy = Inches(2.45) + i * Inches(0.82)
        card(s, rx, yy, rw, Inches(0.7), NAVY if i == 0 else CARD, None if i == 0 else LINE)
        tcol = WHITE if i == 0 else NAVY
        scol = GOLD if i == 0 else GOLD
        mcol = RGBColor(0xC8, 0xD6, 0xE8) if i == 0 else MUTED
        textbox(s, rx + Inches(0.25), yy + Inches(0.08), rw - Inches(0.5), Inches(0.6),
                [{"text": nm + "   ·   " + role, "size": 13.5, "color": tcol, "bold": True, "space_after": 0},
                 {"text": when, "size": 10.5, "color": mcol}])

    # ---------- 8. WORKFLOW DIAGRAM (single combined flow) ----------
    s = add_slide(prs)
    base(s, prs, "Workflow", "One workflow, two deliverables", 8)

    NW, NH = Inches(2.15), Inches(0.98)
    SLOTS = [Inches(0.55), Inches(3.00), Inches(5.45), Inches(7.90), Inches(10.35)]
    ROW1, ROW2 = Inches(2.75), Inches(4.65)
    GREY = RGBColor(0x94, 0xA3, 0xB8)

    def wf_node(x, y, kind, title, sub):
        """kind: step | cp_blue | cp_green | cp_gold | deliverable"""
        if kind == "deliverable":
            shp = card(s, x, y, NW, NH, NAVY, GOLD, line_w=Pt(2.5))
            tcol, scol = WHITE, GOLD
            title = "★  " + title
        elif kind == "step":
            card(s, x, y, NW, NH, CARD, LINE)
            tcol, scol = NAVY, MUTED
        else:
            fill = {"cp_blue": BLUE, "cp_green": GREEN, "cp_gold": GOLD}[kind]
            card(s, x, y, NW, NH, fill, None)
            tcol = NAVY if kind == "cp_gold" else WHITE
            scol = NAVY if kind == "cp_gold" else RGBColor(0xDC, 0xE7, 0xF5)
        lines = [{"text": title, "size": 12.5, "color": tcol, "bold": True,
                  "align": PP_ALIGN.CENTER, "space_after": 1}]
        if sub:
            lines.append({"text": sub, "size": 9.5, "color": scol, "align": PP_ALIGN.CENTER})
        textbox(s, x + Inches(0.06), y, NW - Inches(0.12), NH, lines, anchor=MSO_ANCHOR.MIDDLE)

    def chevron(xc, yc, left=False):
        a = s.shapes.add_shape(MSO_SHAPE.CHEVRON, xc - Inches(0.13), yc - Inches(0.16),
                               Inches(0.26), Inches(0.32))
        _solid(a, GREY)
        a.shadow.inherit = False
        if left:
            a.rotation = 180

    def varrow(xc, y_top, y_bot, shape=MSO_SHAPE.DOWN_ARROW, color=GREY, w=Inches(0.22)):
        a = s.shapes.add_shape(shape, Emu(int(xc) - int(w) // 2), y_top, w, Emu(int(y_bot) - int(y_top)))
        _solid(a, color)
        a.shadow.inherit = False

    # legend
    leg = [("Step", CARD, NAVY), ("Checkpoint", BLUE, WHITE),
           ("Deliverable", NAVY, GOLD), ("Iteration  ↻", WHITE, GOLD)]
    lxp = Inches(0.55)
    for name, fill, _ in leg:
        sw = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lxp, Inches(1.92), Inches(0.26), Inches(0.26))
        sw.fill.solid(); sw.fill.fore_color.rgb = fill
        sw.line.color.rgb = GOLD if name.startswith("Deliverable") else LINE
        sw.line.width = Pt(1.5 if name.startswith("Deliverable") else 1)
        sw.shadow.inherit = False
        tb = textbox(s, lxp + Inches(0.34), Inches(1.9), Inches(2.0), Inches(0.32),
                     [{"text": name, "size": 11, "color": MUTED}], anchor=MSO_ANCHOR.MIDDLE)
        lxp = Emu(int(lxp) + int(Inches(2.15)))

    # ROW 1 (left -> right): metrics -> approval -> DATA HUB
    row1 = [
        (0, "step", "Draft metrics", "one per goal"),
        (1, "cp_blue", "CP1 · Consult", "teams & SMEs"),
        (2, "cp_green", "CP2 · Approve", "metrics · Jacqueline"),
        (3, "step", "Source · validate", "clean data (CP3)"),
        (4, "deliverable", "DATA HUB", "source of truth"),
    ]
    for slot, kind, t, sub in row1:
        wf_node(SLOTS[slot], ROW1, kind, t, sub)
    for slot in range(4):
        xc = Emu(int(SLOTS[slot]) + int(NW) + (int(SLOTS[slot + 1]) - int(SLOTS[slot]) - int(NW)) // 2)
        chevron(xc, Emu(int(ROW1) + int(NH) // 2))

    # WRAP: DATA HUB (row1, slot4) down into Build (row2, slot4)
    hub_cx = Emu(int(SLOTS[4]) + int(NW) // 2)
    varrow(hub_cx, Emu(int(ROW1) + int(NH) + int(Inches(0.04))), Emu(int(ROW2) - int(Inches(0.04))),
           color=GOLD, w=Inches(0.26))

    # ROW 2 (right -> left, snake): Build -> CP4 -> DASHBOARD -> CP5
    row2 = [
        (4, "step", "Build dashboard", "connect to Hub"),
        (3, "cp_gold", "CP4 · Deploy", "IT go/no-go · Nathan"),
        (2, "deliverable", "DASHBOARD", "deployed to execs"),
        (1, "cp_green", "CP5 · Review", "executive feedback"),
    ]
    for slot, kind, t, sub in row2:
        wf_node(SLOTS[slot], ROW2, kind, t, sub)
    for slot in [4, 3, 2]:  # chevrons pointing left between snake nodes
        xc = Emu(int(SLOTS[slot - 1]) + int(NW) + (int(SLOTS[slot]) - int(SLOTS[slot - 1]) - int(NW)) // 2)
        chevron(xc, Emu(int(ROW2) + int(NH) // 2), left=True)

    # ITERATION 1 — big loop: CP5 (row2 slot1) -> bottom channel -> back up into Draft (row1 slot0)
    cp5_cx = Emu(int(SLOTS[1]) + int(NW) // 2)
    chan_y = Inches(6.35)
    varrow(cp5_cx, Emu(int(ROW2) + int(NH) + int(Inches(0.03))), Emu(int(chan_y) + int(Inches(0.05))),
           color=GOLD, w=Inches(0.2))
    channel = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, SLOTS[0], chan_y,
                                 Emu(int(cp5_cx) - int(SLOTS[0]) + int(Inches(0.55))), Inches(0.5))
    _solid(channel, GOLD); channel.shadow.inherit = False
    textbox(s, Emu(int(SLOTS[0]) + int(Inches(0.45))), chan_y, Inches(4.0), Inches(0.5),
            [{"text": "↻  ITERATE — refresh data & metrics", "size": 11, "color": NAVY,
              "bold": True, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    draft_cx = Emu(int(SLOTS[0]) + int(NW) // 2)
    varrow(draft_cx, Emu(int(ROW1) + int(NH) + int(Inches(0.04))), Emu(int(chan_y) - int(Inches(0.02))),
           shape=MSO_SHAPE.UP_ARROW, color=GOLD, w=Inches(0.2))

    # ITERATION 2 — dashboard refine loop: DASHBOARD (slot2) -> Build (slot4), in the gap above row2
    ref_y = Inches(4.26)
    ref_x = Emu(int(SLOTS[2]) + int(NW) // 2)
    ref_w = Emu(int(SLOTS[4]) + int(Inches(0.35)) - int(ref_x))
    ref = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ref_x, ref_y, ref_w, Inches(0.3))
    _solid(ref, NAVY2); ref.shadow.inherit = False
    textbox(s, ref_x, ref_y, ref_w, Inches(0.3),
            [{"text": "↻  refine dashboard", "size": 9.5, "color": WHITE, "bold": True,
              "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)

    # caption
    textbox(s, Inches(0.55), Inches(7.0), Inches(12.2), Inches(0.3),
            [{"text": "A single flow with embedded checkpoints (CP1–CP5). Two iteration loops continuously refine the Data Hub and the dashboard.",
              "size": 11, "color": MUTED, "italic": True, "align": PP_ALIGN.LEFT}])

    # ---------- 9. TIMELINE ----------
    s = add_slide(prs)
    base(s, prs, "Timeline", "A realistic path to end of summer", 9)
    phases = [
        ("IMMEDIATE", "Weeks 1–2", ["Lock scope & draft metrics", "Confirm dataset owners",
                                        "Open IT / deployment talks"], NAVY),
        ("SHORT-TERM", "Weeks 3–5", ["Consult teams (CP1)", "Leadership sign-off (CP2)",
                                          "First validated data in Hub"], BLUE),
        ("MEDIUM-TERM", "Weeks 6–9", ["Populate Hub across domains", "Build dashboard on real data",
                                           "Deployment go/no-go (CP4)"], NAVY2),
        ("END OF SUMMER", "Weeks 10–12", ["Deploy to internal users", "Executive review (CP5)",
                                               "Iterate to refined v1"], GREEN),
    ]
    n = len(phases)
    cw = Inches(2.95)
    gap = Inches(0.18)
    total = n * cw + (n - 1) * gap
    x0 = Emu(int((EMUW - total) / 2))
    # timeline rail
    rail = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x0, Inches(2.35), total, Pt(5))
    _solid(rail, LINE); rail.shadow.inherit = False
    for i, (k, wk, pts, col) in enumerate(phases):
        cx = Emu(int(x0) + i * int(cw + gap))
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + cw/2 - Inches(0.16), Inches(2.18), Inches(0.32), Inches(0.32))
        _solid(dot, col); dot.shadow.inherit = False
        yy = Inches(2.75)
        card(s, cx, yy, cw, Inches(3.4), CARD, LINE)
        head = card(s, cx, yy, cw, Inches(0.95), col, None)
        head.adjustments[0] = 0.07
        textbox(s, cx + Inches(0.2), yy + Inches(0.14), cw - Inches(0.4), Inches(0.8),
                [{"text": k, "size": 14, "color": WHITE, "bold": True, "space_after": 1},
                 {"text": wk, "size": 12, "color": GOLD, "bold": True}])
        for j, p in enumerate(pts):
            py = yy + Inches(1.2) + j * Inches(0.66)
            dot2 = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.25), py + Inches(0.06), Inches(0.13), Inches(0.13))
            _solid(dot2, GOLD)
            textbox(s, cx + Inches(0.5), py, cw - Inches(0.7), Inches(0.6),
                    [{"text": p, "size": 12, "color": INK}])
    textbox(s, Inches(0.55), Inches(6.45), Inches(12.2), Inches(0.4),
            [{"text": "Checkpoints gate progress:  CP1 → CP2 → CP3 → CP4 → CP5.  Deployment depends on IT provisioning (start early).",
              "size": 12.5, "color": MUTED, "italic": True, "align": PP_ALIGN.CENTER}])

    # ---------- 10. RISKS ----------
    s = add_slide(prs)
    base(s, prs, "Risks", "Main bottlenecks & how we manage them", 10)
    risks = [
        ("Finding reliable data", "Prefer public sources; allow internal/proxy metrics; flag gaps openly"),
        ("Validating credibility", "SME validation before any figure enters the Hub"),
        ("Timely team input", "Owner packets + biweekly check-ins; escalate when blocked"),
        ("Leadership approval", "Early consultation so sign-off is a confirmation, not a debate"),
        ("Clean, reliable Hub", "Controlled vocabularies, ownership, periodic clean-up"),
        ("Deploying the dashboard", "Engage government IT early, in parallel (key bottleneck)"),
    ]
    cw, ch = Inches(3.95), Inches(1.9)
    gx, gy = Inches(0.19), Inches(0.25)
    x0, y0 = Inches(0.55), Inches(2.25)
    for i, (t, d) in enumerate(risks):
        r, c = divmod(i, 3)
        x = x0 + c * (cw + gx)
        yy = y0 + r * (ch + gy)
        card(s, x, yy, cw, ch, CARD, LINE)
        accent_bar(s, x + Inches(0.28), yy + Inches(0.28), Inches(0.5), GOLD)
        textbox(s, x + Inches(0.28), yy + Inches(0.5), cw - Inches(0.56), Inches(1.3),
                [{"text": t, "size": 15, "color": NAVY, "bold": True, "space_after": 5},
                 {"text": d, "size": 12, "color": MUTED}])

    # ---------- 11. NEXT STEPS / CLOSING ----------
    s = add_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMUW, EMUH)
    bg.fill.gradient()
    bg.fill.gradient_stops[0].color.rgb = NAVY
    bg.fill.gradient_stops[1].color.rgb = NAVY2
    try:
        bg.fill.gradient_angle = 30.0
    except Exception:
        pass
    _noline(bg); bg.shadow.inherit = False
    accent_bar(s, Inches(0.9), Inches(0.85), Inches(1.2), GOLD, Pt(6))
    textbox(s, Inches(0.9), Inches(1.05), Inches(11), Inches(0.9),
            [{"text": "Immediate next steps", "size": 36, "color": WHITE, "bold": True, "font": FONT_L}])
    steps = [
        ("Confirm dataset owners", "From the Workflow & Ownership Registry"),
        ("Draft metrics per goal", "Mehdi + Daniel, using the standard template"),
        ("Open deployment talks", "Nathan / government IT — start in parallel now"),
        ("Schedule consultation (CP1)", "Internal teams & SMEs confirm metrics"),
        ("Book leadership sign-off (CP2)", "Lock the approved metric set with Jacqueline"),
        ("Stand up the cadence", "Weekly internal sync + biweekly owner check-ins"),
    ]
    cw, ch = Inches(5.7), Inches(1.25)
    gx, gy = Inches(0.5), Inches(0.3)
    x0, y0 = Inches(0.9), Inches(2.4)
    for i, (t, d) in enumerate(steps):
        r, c = divmod(i, 2)
        x = x0 + c * (cw + gx)
        yy = y0 + r * (ch + gy)
        num = s.shapes.add_shape(MSO_SHAPE.OVAL, x, yy + Inches(0.1), Inches(0.7), Inches(0.7))
        _solid(num, GOLD)
        textbox(s, x, yy + Inches(0.1), Inches(0.7), Inches(0.7),
                [{"text": str(i+1), "size": 24, "color": NAVY, "bold": True, "align": PP_ALIGN.CENTER}],
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, x + Inches(0.95), yy, cw - Inches(1.0), ch,
                [{"text": t, "size": 17, "color": WHITE, "bold": True, "space_after": 2},
                 {"text": d, "size": 12.5, "color": RGBColor(0xC0, 0xCE, 0xE0)}],
                anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4),
            [{"text": "Target: process streamlined and dashboard deployed by end of summer.",
              "size": 14, "color": GOLD, "bold": True}])

    prs.save(dst)
    print(f"Saved {dst}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build(sys.argv[1])
